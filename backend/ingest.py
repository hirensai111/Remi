from pathlib import Path

import fitz


def ingest_file(path: str) -> tuple[str, int, str]:
    """
    Extract text from a file.

    Returns:
        (extracted_text, unit_count, file_type)
        unit_count = pages (pdf), slides (pptx), paragraphs (docx), lines (txt/md)
    """
    ext = Path(path).suffix.lower()

    if ext == ".pdf":
        return _ingest_pdf(path)
    elif ext == ".docx":
        return _ingest_docx(path)
    elif ext == ".pptx":
        return _ingest_pptx(path)
    elif ext in (".txt", ".md"):
        return _ingest_text(path)
    else:
        raise ValueError(
            f"Unsupported file type: '{ext}'. Supported: .pdf, .docx, .pptx, .txt, .md"
        )


def _ingest_pdf(path: str) -> tuple[str, int, str]:
    doc = fitz.open(path)
    page_count = len(doc)
    pages: list[str] = []
    for i in range(page_count):
        page = doc.load_page(i)
        text = page.get_text().strip()
        if text:
            pages.append(f"[Page {i + 1}]\n{text}")
    doc.close()
    return "\n\n".join(pages), page_count, "pdf"


def _ingest_docx(path: str) -> tuple[str, int, str]:
    from docx import Document

    doc = Document(path)
    parts: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                parts.append(" | ".join(row_text))

    return "\n\n".join(parts), len(parts), "docx"


def _ingest_pptx(path: str) -> tuple[str, int, str]:
    from pptx import Presentation

    prs = Presentation(path)
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides), len(prs.slides), "pptx"


def _ingest_text(path: str) -> tuple[str, int, str]:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    lines = text.splitlines()
    return text, len(lines), "text"
